"""Generate minimal test fixtures (.xlsx, .docx, .pptx) for smoke-testing."""

import os, sys
sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))

import openpyxl
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt

OUT = os.path.join(os.path.dirname(__file__), "fixtures")
os.makedirs(OUT, exist_ok=True)


def make_excel():
    # --- Version 1 ---
    wb1 = openpyxl.Workbook()
    ws1 = wb1.active
    ws1.title = "Revenue"
    ws1.append(["Region", "Q1", "Q2", "Q3", "Q4"])
    ws1.append(["North", "120000", "135000", "128000", "142000"])
    ws1.append(["South", "98000", "102000", "97000", "105000"])
    ws1.append(["Contact: John Smith", "Email: john@example.com", "SSN: 123-45-6789", "", ""])
    wb1.save(os.path.join(OUT, "report_v1.xlsx"))

    # --- Version 2 ---
    wb2 = openpyxl.Workbook()
    ws2 = wb2.active
    ws2.title = "Revenue"
    ws2.append(["Region", "Q1", "Q2", "Q3", "Q4"])
    ws2.append(["North", "120000", "138000", "128000", "145000"])
    ws2.append(["South", "98000", "102000", "99000", "105000"])
    ws2.append(["East", "75000", "78000", "81000", "83000"])
    ws2.append(["Contact: Jane Doe", "Email: jane@example.com", "SSN: 987-65-4321", "", ""])
    wb2.save(os.path.join(OUT, "report_v2.xlsx"))


def make_word():
    # --- Version 1 ---
    d1 = Document()
    d1.add_heading("Quarterly Financial Summary", level=1)
    d1.add_paragraph("Prepared by: Alice Johnson")
    d1.add_paragraph("Total revenue for the period was $503,000.")
    d1.add_paragraph("Operating expenses totalled $312,000.")
    d1.add_paragraph("Net income: $191,000.")
    d1.add_paragraph("Contact: 555-123-4567 | alice@corp.com")

    t1 = d1.add_table(rows=3, cols=3)
    t1.rows[0].cells[0].text = "Category"
    t1.rows[0].cells[1].text = "Budget"
    t1.rows[0].cells[2].text = "Actual"
    t1.rows[1].cells[0].text = "Marketing"
    t1.rows[1].cells[1].text = "50000"
    t1.rows[1].cells[2].text = "48000"
    t1.rows[2].cells[0].text = "Engineering"
    t1.rows[2].cells[1].text = "120000"
    t1.rows[2].cells[2].text = "118000"

    d1.save(os.path.join(OUT, "summary_v1.docx"))

    # --- Version 2 ---
    d2 = Document()
    d2.add_heading("Quarterly Financial Summary", level=1)
    d2.add_paragraph("Prepared by: Bob Williams")
    d2.add_paragraph("Total revenue for the period was $527,000.")
    d2.add_paragraph("Operating expenses totalled $318,000.")
    d2.add_paragraph("Net income: $209,000.")
    d2.add_paragraph("Contact: 555-987-6543 | bob@corp.com")

    t2 = d2.add_table(rows=4, cols=3)
    t2.rows[0].cells[0].text = "Category"
    t2.rows[0].cells[1].text = "Budget"
    t2.rows[0].cells[2].text = "Actual"
    t2.rows[1].cells[0].text = "Marketing"
    t2.rows[1].cells[1].text = "55000"
    t2.rows[1].cells[2].text = "53000"
    t2.rows[2].cells[0].text = "Engineering"
    t2.rows[2].cells[1].text = "120000"
    t2.rows[2].cells[2].text = "122000"
    t2.rows[3].cells[0].text = "Sales"
    t2.rows[3].cells[1].text = "80000"
    t2.rows[3].cells[2].text = "79000"

    d2.save(os.path.join(OUT, "summary_v2.docx"))


def make_pptx():
    prs = Presentation()
    layout = prs.slide_layouts[1]

    slide1 = prs.slides.add_slide(layout)
    slide1.shapes.title.text = "Q3 2025 Business Review"
    slide1.placeholders[1].text = (
        "Prepared by: Carol Davis\n"
        "Period: Q3 2025 (Jul–Sep)\n"
        "Total Revenue: $4.2M\n"
        "Operating Margin: 18.5%\n"
        "Headcount: 342\n"
        "Customer NPS: 72"
    )

    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    slide2.shapes.title.text = "Regional Performance"
    rows, cols = 5, 4
    left = Inches(0.5)
    top = Inches(1.8)
    width = Inches(9.0)
    height = Inches(2.5)
    table = slide2.shapes.add_table(rows, cols, left, top, width, height).table
    headers = ["Region", "Revenue", "Growth %", "Headcount"]
    for ci, h in enumerate(headers):
        table.cell(0, ci).text = h
    data = [
        ["North America", "$2.1M", "12%", "180"],
        ["Europe", "$1.3M", "8%", "95"],
        ["Asia-Pacific", "$0.6M", "22%", "45"],
        ["Latin America", "$0.2M", "5%", "22"],
    ]
    for ri, row in enumerate(data, start=1):
        for ci, val in enumerate(row):
            table.cell(ri, ci).text = val

    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Key Metrics"
    slide3.placeholders[1].text = (
        "Customer Acquisition Cost: $145\n"
        "Lifetime Value: $2,800\n"
        "LTV/CAC Ratio: 19.3:1\n"
        "Monthly Churn Rate: 1.2%\n"
        "ARR: $16.8M\n"
        "Burn Rate: $380K/month\n"
        "Runway: 18 months\n"
        "Contact: carol@corp.com | 555-234-5678"
    )

    prs.save(os.path.join(OUT, "deck_q3.pptx"))


if __name__ == "__main__":
    make_excel()
    make_word()
    make_pptx()
    print("Fixtures created in", OUT)
